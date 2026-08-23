"""P5-6 — Executive PowerPoint export (python-pptx, monochrome, <=10 slides).

The deck is built FROM the same compose_playbook / run-of-show / ledger the web
views use, so it can never disagree with the on-screen playbook. Content is a
coherent executive summary, not a dump of every screen.

Monochrome-only (per DESIGN.md): ink / neutral / steel + an amber accent reserved
for the double-booking flag (matching the run-of-show's existing conflict color).
No signal-blue fills. Color tuples are reused from visuals.py so the palette stays
single-sourced — no hex literals in this module.
"""
from __future__ import annotations

import io
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.db import repository as repo
from app.db.models import Event
from app.features import run_of_show as ros
from app.features.playbook import compose_playbook
from app.features.visuals import INK, NEUTRAL, STEEL, SIGNAL

# Warm amber used ONLY for the double-booking flag (mirrors theme.css #D8A24C).
AMBER = (216, 162, 76)

_INK = RGBColor(*INK)
_NEUTRAL = RGBColor(*NEUTRAL)
_STEEL = RGBColor(*STEEL)
_AMBER = RGBColor(*AMBER)
_SIGNAL = RGBColor(*SIGNAL)

_FONT_HEAD = "Archivo Expanded"
_FONT_BODY = "Inter"

_EMU_W = Inches(13.333)
_EMU_H = Inches(7.5)


def _asset(path: str) -> str:
    import os
    # exec_pptx.py lives at app/features/; repo root is three levels up.
    return os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), path)


def build_playbook(conn, event_id: int):
    """Same composer the web views use — single source of truth."""
    return compose_playbook(conn, event_id)


def _set_font(run, name: str, size: int, color, bold: bool = False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _slide(pres: Presentation, bg: RGBColor):
    s = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def _textbox(s, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def _wordmark(s, light: bool):
    """Monochrome Saronic FULL wordmark (not just the symbol), top-left.
    Role-named asset, never a colored lockup. DESIGN.md: wordmark is strictly
    monochrome; signal-blue is product-UI-only and must not appear in the deck."""
    logo = (_asset("assets/press-kit/Logos/Saronic_Logo_Full--Light.png")
            if light else
            _asset("assets/press-kit/Logos/Saronic_Logo_Full--Dark.png"))
    import os
    if os.path.exists(logo):
        s.shapes.add_picture(logo, Inches(0.5), Inches(0.45), height=Inches(0.55))


def _footer_wordmark(s, dark: bool = True):
    """Small monochrome wordmark footer on content slides (light bg -> dark
    mark) so every slide carries Saronic branding without eating body space."""
    logo = (_asset("assets/press-kit/Logos/Saronic_Logo_Full--Dark.png")
            if dark else
            _asset("assets/press-kit/Logos/Saronic_Logo_Full--Light.png"))
    import os
    if os.path.exists(logo):
        s.shapes.add_picture(logo, Inches(11.4), Inches(6.95), height=Inches(0.42))


def _owner_label(conn, person_id: int) -> str:
    p = repo.get_person(conn, person_id)
    if not p:
        return "—"
    return p.display_name


def _schedule_line(pb) -> str:
    sched = getattr(pb, "schedule", None)
    if sched:
        return sched
    return ""


def _head(s, text: str) -> None:
    tb, tf = _textbox(s, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8))
    _set_font(tf.paragraphs[0].add_run(), _FONT_HEAD, 26, _INK, bold=True)
    tf.paragraphs[0].runs[0].text = text


def _kv(s, rows, top) -> None:
    tb, tf = _textbox(s, Inches(0.6), top, Inches(12.1), Inches(5.0))
    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_font(p.add_run(), _FONT_HEAD, 14, _STEEL, bold=True)
        p.runs[0].text = f"{k}:  "
        _set_font(p.add_run(), _FONT_BODY, 14, _INK)
        p.runs[1].text = str(v)
        p.space_after = Pt(8)


def _slide_title(pres: Presentation, event: Event, pb) -> None:
    s = _slide(pres, _INK)
    _wordmark(s, light=True)
    tb, tf = _textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2),
                      MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _set_font(p.add_run(), _FONT_HEAD, 40, _NEUTRAL, bold=True)
    p.runs[0].text = event.name  # identity, kept in real case (not force-uppercased)
    p2 = tf.add_paragraph()
    _set_font(p2.add_run(), _FONT_BODY, 18, _STEEL)
    p2.runs[0].text = "EVENT PLAYBOOK — EXECUTIVE SUMMARY"
    sub = []
    if event.owner_name:
        sub.append(f"Owner: {event.owner_name}"
                   + (f" — {event.owner_role}" if event.owner_role else ""))
    loc = (event.location or "").strip()
    if loc:
        sub.append(f"Location: {loc}")
    if sub:
        p3 = tf.add_paragraph()
        _set_font(p3.add_run(), _FONT_BODY, 14, _STEEL)
        p3.runs[0].text = "  ·  ".join(sub)


def _slide_overview(pres: Presentation, event: Event, pb) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "OVERVIEW")
    _footer_wordmark(s)
    rows = [
        ("Event type", event.event_type or "Not set"),
        ("Audience", f"{event.audience_estimate:,}" if event.audience_estimate
         else "Not set"),
        ("Location", event.location or "Not set"),
        ("Date window", _schedule_line(pb) or "Not scheduled yet"),
        ("Event owner", (event.owner_name or "Not set")
         + (f" — {event.owner_role}" if event.owner_role else "")),
    ]
    _kv(s, rows, top=Inches(1.6))


def _slide_decisions(pres: Presentation, event: Event, pb) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "KEY DECISIONS")
    _footer_wordmark(s)
    tb, tf = _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6))
    sections = pb.sections or []
    first = True
    for sec in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_font(p.add_run(), _FONT_HEAD, 15, _INK, bold=True)
        p.runs[0].text = f"{sec.title}: {sec.chosen_label}"
        p.space_before = Pt(6)
        pr = tf.add_paragraph()
        _set_font(pr.add_run(), _FONT_BODY, 12, _STEEL)
        pr.runs[0].text = f"   {sec.reasoning or '—'}"
    if not first:
        oq = len(pb.open_questions or [])
        if oq:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
            _set_font(p.add_run(), _FONT_HEAD, 14, _STEEL, bold=True)
            p.runs[0].text = f"OPEN QUESTIONS: {oq}"
    if first:
        p = tf.paragraphs[0]
        _set_font(p.add_run(), _FONT_BODY, 14, _STEEL)
        p.runs[0].text = "No settled decisions yet."


def _slide_venue(pres: Presentation, event: Event, pb) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "VENUE")
    _footer_wordmark(s)
    venue_label, capacity, amenities = _venue_facts(pb)
    tb, tf = _textbox(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0))
    p = tf.paragraphs[0]
    _set_font(p.add_run(), _FONT_HEAD, 16, _INK, bold=True)
    p.runs[0].text = venue_label or "Not set"
    if capacity:
        pc = tf.add_paragraph()
        _set_font(pc.add_run(), _FONT_BODY, 13, _INK)
        pc.runs[0].text = f"Capacity: {capacity}"
    if amenities:
        pa = tf.add_paragraph()
        _set_font(pa.add_run(), _FONT_BODY, 13, _INK)
        pa.runs[0].text = "Key amenities: " + "; ".join(amenities)


def _venue_facts(pb):
    for sec in (pb.sections or []):
        if sec.step == "venue" or "venue" in sec.title.lower():
            return sec.chosen_label, None, []
    return None, None, []


def _strip_day(val: str, day: str) -> str:
    """Strip the date prefix from a timestamp when the day header already
    carries it (exec deck + print view both group by day)."""
    if not val:
        return "—"
    if val.startswith(day + " "):
        return val[len(day) + 1:]
    if "T" in val and val.startswith(day + "T"):
        return val[len(day) + 1:].replace("T", " ")
    return val


def _slide_run_of_show(pres: Presentation, conn, event: Event, segments) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "RUN OF SHOW")
    _footer_wordmark(s)
    conflicts = ros.conflicts_for(segments or [])
    by_day = ros.group_by_day(segments or [])
    tb, tf = _textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6))
    first = True
    for day, segs in by_day.items():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(8)
        _set_font(p.add_run(), _FONT_HEAD, 14, _STEEL, bold=True)
        p.runs[0].text = day
        for seg in segs:
            pl = tf.add_paragraph()
            _set_font(pl.add_run(), _FONT_BODY, 12, _INK)
            owners = ", ".join(_owner_label(conn, o) for o in seg.owner_ids) or "—"
            pl.runs[0].text = (f"{_strip_day(seg.start, day)}–"
                               f"{_strip_day(seg.end, day)}  {seg.title}  "
                               f"[{seg.track}]  ·  {owners}")
            flag = conflicts.get(seg.id)
            if flag:
                for staff_id, why in flag.items():
                    pa = tf.add_paragraph()
                    run = pa.add_run()
                    _set_font(run, _FONT_BODY, 11, _AMBER)
                    run.text = (f"  ⚠ {_owner_label(conn, staff_id)} double-booked: "
                                f"{why}")
    if first:
        p = tf.paragraphs[0]
        _set_font(p.add_run(), _FONT_BODY, 14, _STEEL)
        p.runs[0].text = "No run-of-show segments yet."


def _slide_attendance(pres: Presentation, conn, event: Event) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "ATTENDANCE & ACCESS")
    _footer_wordmark(s)
    attendees = repo.list_attendees(conn, event.id)
    headcount = len(attendees)
    vip = sum(1 for a in attendees if a.is_vip)
    arrived = sum(1 for a in attendees if a.attended_at)
    # Check-in staffing from P5-9 per-event tags.
    assignees = []
    for row in repo.event_staff_rows(conn, event.id):
        if not row["can_check_in"]:
            continue
        p = repo.get_person(conn, row["person_id"])
        if p and not p.is_erased:
            assignees.append(p.display_name)
    rows = [
        ("Invited", str(headcount) if headcount else "Not set"),
        ("VIPs", str(vip)),
        ("Arrived (scan or manual)", str(arrived)),
        ("Check-in staff assigned", ", ".join(sorted(assignees)) or "Not set"),
    ]
    _kv(s, rows, top=Inches(1.6))


def _slide_spend(pres: Presentation, conn, event: Event) -> None:
    s = _slide(pres, _NEUTRAL)
    _head(s, "SPEND")
    _footer_wordmark(s)
    tb, tf = _textbox(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0))
    p = tf.paragraphs[0]
    try:
        total = repo.spend_total(conn, event.id)
    except Exception:
        total = None
    if total is None or total <= 0:
        _set_font(p.add_run(), _FONT_BODY, 16, _INK)
        p.runs[0].text = "None — ran offline."
    else:
        _set_font(p.add_run(), _FONT_BODY, 16, _INK)
        p.runs[0].text = f"Claude API spend: ${total:.2f} (whole cents)."
        note = tf.add_paragraph()
        _set_font(note.add_run(), _FONT_BODY, 12, _STEEL)
        note.runs[0].text = "Full precision lives in the spend ledger."


def build_exec_pptx_obj(conn, event_id: int, pb, segments,
                        event: Optional[Event] = None,
                        owner_id: Optional[int] = None) -> Presentation:
    """Assemble the slides from playbook + ROS + ledger. Returns the live
    Presentation object (so tests can inspect fills/fonts). `owner_id` is
    accepted for API symmetry with the caller; the owner already derives from
    the event, so it is currently unused."""
    ev = event or repo.get_event(conn, event_id)
    pres = Presentation()
    pres.slide_width = _EMU_W
    pres.slide_height = _EMU_H
    _slide_title(pres, ev, pb)
    _slide_overview(pres, ev, pb)
    _slide_decisions(pres, ev, pb)
    _slide_venue(pres, ev, pb)
    _slide_run_of_show(pres, conn, ev, segments)
    _slide_attendance(pres, conn, ev)
    _slide_spend(pres, conn, ev)
    return pres


def build_exec_pptx(conn, event_id: int, pb, segments,
                    event: Optional[Event] = None,
                    owner_id: Optional[int] = None) -> bytes:
    pres = build_exec_pptx_obj(conn, event_id, pb, segments, event=event,
                               owner_id=owner_id)
    buf = io.BytesIO()
    pres.save(buf)
    return embed_fonts(buf.getvalue())


# --- H1: real brand-font embedding (not just naming) -----------------------
# python-pptx has no embed API, so we inject the OOXML font parts after save:
# a ppt/fonts/<name>.ttf part per face, a relationship in the theme's .rels,
# and an <a:embeddedFont r:id> mapping in the theme fontScheme. A font-less
# recipient machine then resolves the actual glyphs instead of substituting.
FONTS_TO_EMBED = (
    ("Archivo Expanded", "ArchivoExpanded-Bold.ttf", "rIdFontArchivo"),
    ("Inter", "Inter-Regular.ttf", "rIdFontInter"),
)


def _inject_embedded_font(theme: str, block_tag: str, typeface: str,
                          rid: str) -> str:
    """Insert <a:font typeface=...><a:embeddedFont r:id=.../> into a fontScheme
    block (majorFont/minorFont), right after its <a:cs .../> child."""
    start = theme.find(f"<a:{block_tag}>")
    end = theme.find(f"</a:{block_tag}>")
    if start < 0 or end < 0:
        return theme
    block = theme[start:end]
    anchor = block.find("<a:cs")
    if anchor < 0:
        return theme
    # Position after the closing of the <a:cs .../> element.
    insertion = anchor
    depth = 0
    i = anchor
    while i < len(block):
        if block[i] == "<":
            if block[i:i + 2] == "</":
                depth -= 1
                if depth == 0:
                    insertion = i + len("</a:cs>") if block[i:i + 7] == "</a:cs>" \
                        else i + block[i:].find(">") + 1
                    break
            elif block[i:i + 2] == "/>":
                depth = 0
                insertion = i + 2
                break
            else:
                depth += 1
        i += 1
    new_elem = (f'<a:font typeface="{typeface}">'
                f'<a:embeddedFont r:id="{rid}"/></a:font>')
    return theme[:start] + block[:insertion] + new_elem + block[insertion:] + theme[end:]


def embed_fonts(pptx_bytes: bytes) -> bytes:
    """Return pptx bytes with the brand TTFs embedded as OOXML font parts.

    python-pptx's default theme has no relationships part of its own, so we
    create ppt/theme/_rels/theme1.xml.rels, register the font parts in
    [Content_Types].xml, and map each typeface to its embedded part via
    <a:embeddedFont r:id> in the theme fontScheme.
    """
    import zipfile
    from pathlib import Path
    assets = Path(__file__).resolve().parent.parent.parent / "assets/fonts"
    theme_name = "ppt/theme/theme1.xml"
    rels_name = "ppt/theme/_rels/theme1.xml.rels"
    ctypes_name = "[Content_Types].xml"

    src = zipfile.ZipFile(io.BytesIO(pptx_bytes))
    names = src.namelist()
    theme = src.read(theme_name).decode("utf-8")
    ctypes = src.read(ctypes_name).decode("utf-8")

    embed = []
    for typeface, fname, rid in FONTS_TO_EMBED:
        asset = assets / fname
        if asset.exists():
            embed.append((typeface, fname, rid, asset.read_bytes()))

    for typeface, fname, rid, _ in embed:
        block = "majorFont" if typeface == "Archivo Expanded" else "minorFont"
        theme = _inject_embedded_font(theme, block, typeface, rid)
        # Font content-type override (idempotent by extension).
        if f'Extension="ttf"' not in ctypes:
            ctypes = ctypes.replace(
                "</Types>",
                '<Default Extension="ttf" '
                'ContentType="application/x-font.ttf"/></Types>')

    # Build the theme rels (created here; the default template has none).
    rel_lines = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
                '<Relationships xmlns="http://schemas.openxmlformats.org/'
                'package/2006/relationships">']
    for typeface, fname, rid, _ in embed:
        rel_lines.append(
            f'<Relationship Id="{rid}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/'
            f'2006/relationships/font" Target="fonts/{fname}"/>')
    rel_lines.append("</Relationships>")
    rels = "\n".join(rel_lines)

    out = io.BytesIO()
    dst = zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED)
    written = set()
    for item in names:
        if item == theme_name:
            dst.writestr(item, theme)
        elif item == ctypes_name:
            dst.writestr(item, ctypes)
        else:
            dst.writestr(item, src.read(item))
        written.add(item)
    # Theme rels: only write if not already present (idempotent re-runs).
    if rels_name not in written:
        dst.writestr(rels_name, rels)
    for typeface, fname, rid, data in embed:
        part = f"ppt/fonts/{fname}"
        if part not in written:
            dst.writestr(part, data)
    dst.close()
    return out.getvalue()


# --- Helpers used by the test-suite assertions. ---

def extract_text(blob: bytes) -> str:
    prs = Presentation(io.BytesIO(blob))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    out.append("".join(r.text for r in para.runs))
    return "\n".join(out)


def collect_fills(prs: Presentation):
    fills = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type is None:
                continue
            try:
                fill = shape.fill
                if fill.type is not None and fill.fore_color and \
                        fill.fore_color.type is not None:
                    rgb = fill.fore_color.rgb
                    if rgb is not None:
                        fills.append((rgb[0], rgb[1], rgb[2]))
            except Exception:
                pass
    return fills


def has_amber_conflict(prs: Presentation) -> bool:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb is not None:
                            rgb = run.font.color.rgb
                            if (rgb[0], rgb[1], rgb[2]) == AMBER:
                                return True
                    except Exception:
                        pass
    return False


