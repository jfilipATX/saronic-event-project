"""P7-5 — PPTX branding + tighter layout.

The exec deck now carries the monochrome FULL wordmark (not just the symbol):
on the title slide (light mark, dark bg) and as a footer on every content
slide. Layout frames use the wider 12.1" content width (0.6" margins). Brand
rule preserved: no signal-blue in the deck.
"""
from __future__ import annotations

import os
import sqlite3
import tempfile

from fastapi.testclient import TestClient

from app.db import repository as repo
from app.db.models import Event, Person, Segment
from app.features import exec_pptx as ex
from app.main import create_app


def _client_db():
    d = tempfile.mkdtemp()
    os.environ["DB_PATH"] = os.path.join(d, "t.db")
    return TestClient(create_app())


def _make(client):
    r = client.post("/events", data={"name": "Fleet Week", "city": "SF",
                                     "state": "CA", "country": "US",
                                     "owner_name": "Sam", "owner_role": "Lead"},
                    follow_redirects=False)
    return int(r.headers["location"].rstrip("/").split("/")[2])


def _direct():
    c = sqlite3.connect(os.environ["DB_PATH"], isolation_level=None)
    c.row_factory = sqlite3.Row
    return c


def _seed(c, eid):
    pid = repo.add_person(c, Person(name="Sam", role="Ops"))
    repo.assign_staff(c, eid, pid, role="Ops")
    repo.add_segment(c, Segment(event_id=eid, title="Expo",
        start="2026-05-08 09:00", end="2026-05-08 18:00", kind="floor",
        track="Floor", owner_ids=[pid]))


def test_deck_carries_full_wordmark_on_title_and_footers():
    client = _client_db()
    eid = _make(client)
    c = _direct()
    _seed(c, eid)
    # Build directly via the feature (single source of truth).
    pb = ex.build_playbook(c, eid)
    segs = repo.list_segments(c, eid)
    blob = ex.build_exec_pptx(c, eid, pb, segs)
    # Use python-pptx to count pictures.
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(blob))
    pics = [sh for s in prs.slides for sh in s.shapes if sh.shape_type == 13]
    # Title full wordmark + 6 content-slide footers = 7 pictures minimum.
    assert len(pics) >= 7


def test_deck_is_monochrome_no_signal_blue():
    client = _client_db()
    eid = _make(client)
    c = _direct()
    _seed(c, eid)
    pb = ex.build_playbook(c, eid)
    segs = repo.list_segments(c, eid)
    blob = ex.build_exec_pptx(c, eid, pb, segs)
    import io
    import pptx
    prs = pptx.Presentation(io.BytesIO(blob))
    fills = ex.collect_fills(prs)
    SIGNAL = (76, 159, 216)
    assert SIGNAL not in fills
